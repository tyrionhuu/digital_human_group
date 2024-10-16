import os
from PIL import Image
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO


def pdf_to_pptx(pdf_file):
    print(f"\nConverting file: {pdf_file}\n")

    # Prep presentation
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    # Create base name for output
    base_name = os.path.splitext(pdf_file)[0]

    # Convert PDF to list of images
    print("Starting conversion...")
    slideimgs = convert_from_path(pdf_file, 300, fmt='ppm', thread_count=2)
    print("...complete.\n")

    # Loop over slides
    for i, slideimg in enumerate(slideimgs):
        if i % 10 == 0:
            print(f"Saving slide: {i}")

        # Save image to BytesIO
        imagefile = BytesIO()
        slideimg.save(imagefile, format='tiff')
        imagefile.seek(0)
        width, height = slideimg.size

        # Set slide dimensions
        prs.slide_height = height * 9525
        prs.slide_width = width * 9525

        # Add slide to presentation
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(imagefile, 0, 0, width=width * 9525, height=height * 9525)

    # Save PowerPoint file
    output_file = base_name + '.pptx'
    print(f"\nSaving file: {output_file}")
    prs.save(output_file)
    print("Conversion complete. :)\n")

# Example usage
pdf_to_pptx('../test/Jacob_Devlin_BERT.pdf')
