from pptx import Presentation
import zipfile
import os

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return text

def pptx_to_xml(pptx_path, output_folder):
    # Ensure the output folder exists
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # Extract the .pptx (zip) file to the output folder
    with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
        zip_ref.extractall(output_folder)

    print(f"PPTX content extracted to {output_folder}")


# Usage example
pptx_path = "../test/test.pptx"
output_folder = "extracted_content"
pptx_to_xml(pptx_path, output_folder)