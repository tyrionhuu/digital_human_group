from pptx import Presentation
import xml.etree.ElementTree as ET
import os


def xml2ppt(xml_file, ppt_file):
    # Load the XML file
    tree = ET.parse(xml_file)
    root = tree.getroot()

    # Create a new PowerPoint presentation
    prs = Presentation()

    # Loop through each slide element in the XML
    for slide_element in root.findall('slide'):
        slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
        slide = prs.slides.add_slide(slide_layout)

        # Get the slide title from the XML
        title_element = slide_element.find('title')
        title = title_element.text if title_element is not None else ''

        # Add the title to the slide (if there's a title and title placeholder exists)
        if title and slide.shapes.title:
            slide.shapes.title.text = title

        # Add content to the slide
        content_element = slide_element.find('content')
        if content_element is not None:
            content = content_element.text
            if content:
                # Add content to the slide (if content placeholder exists)
                if slide.shapes.placeholders:
                    for shape in slide.shapes.placeholders:
                        if shape.placeholder_format.idx == 1:
                            shape.text = content

    # Save the PowerPoint file
    prs.save(ppt_file)
    print(f'Successfully created PowerPoint: {ppt_file}')


# Usage example
xml_file_path = '../test/test.xml'  # Modify the path as needed
ppt_file_path = '../test/generated_ppt.pptx'  # Modify the path as needed

if os.path.exists(xml_file_path):
    xml2ppt(xml_file_path, ppt_file_path)
else:
    print(f"XML file not found: {xml_file_path}")
