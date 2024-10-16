from pptx import Presentation
import xml.etree.ElementTree as ET
import re
from xml.sax.saxutils import escape
from xml.dom import minidom
import os
import ollama
IMAGE_DIR = "../test/images"
def call_llava(image_path):
    """
    Calls the LLaVA model via Ollama API to describe images, charts, or tables.
    """
    try:
        # Call LLaVA through Ollama to get the description
        res = ollama.chat(
            model="llava",
            messages=[
                {
                    'role': 'user',
                    'content': 'Describe this image:',
                    'images': [image_path]
                }
            ]
        )
        # Extract the description from the response
        description = res['message']['content']
        return description
    except Exception as e:
        print(f"Error calling LLaVA for {image_path}: {e}")
        return f"Error describing {image_path}"

def filter_text(text):
    # Replace control characters and invalid characters with a space or remove them
    cleaned_text = re.sub(r'[\x00-\x1F\x7F]', '', text)  # Remove control characters
    return cleaned_text

def save_image_from_shape(shape, slide_num, shape_num, output_dir=IMAGE_DIR):
    """Save an image from a shape to a file."""
    if not hasattr(shape, 'image'):
        return None

    image = shape.image
    image_filename = f"slide_{slide_num}_shape_{shape_num}.{image.ext}"
    image_path = os.path.join(output_dir, image_filename)

    # Save the image blob to a file
    with open(image_path, 'wb') as f:
        f.write(image.blob)

    return image_path
def ppt2xml(ppt_file, xml_file):
    # Load the PowerPoint presentation
    prs = Presentation(ppt_file)

    # Create the root element for the XML
    root = ET.Element('presentation')

    # Loop through each slide in the presentation
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_element = ET.SubElement(root, 'slide', number=str(slide_num))
        # Get the title of the slide
        title = ''
        if slide.shapes.title:
            title = slide.shapes.title.text

        # Create an element for the title (if not empty)
        title = filter_text(title)
        if title:
            title_element = ET.SubElement(slide_element, 'title')
            title_element.text = escape(title)

        # Loop through each shape in the slide to capture bullet points, content, and images
        for shape_num, shape in enumerate(slide.shapes, start=1):
            print(shape)
            if not hasattr(shape, 'has_text_frame') and not hasattr(shape, 'has_image'):
                continue  # Skip shapes without text or images

            # Handle text shapes
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    para_text = filter_text(para.text)  # Filter the paragraph text

                    # Skip empty paragraphs
                    if not para_text:
                        continue

                    # Create a new paragraph element
                    content_element = ET.SubElement(slide_element, 'paragraph')

                    # Check for bullet points
                    if para.level > 0:
                        content_element.set('type', 'bullet')
                        content_element.set('level', str(para.level))
                    else:
                        content_element.set('type', 'normal')

                    # Assign the filtered text to the paragraph element
                    content_element.text = escape(para_text)

            # Handle image shapes
            if hasattr(shape, 'image'):
                image_path = save_image_from_shape(shape, slide_num, shape_num)
                if image_path:
                    # Call LLaVA to get the description of the image
                    description = call_llava(image_path)
                    image_element = ET.SubElement(slide_element, 'image', number=str(shape_num))
                    image_element.text = escape(description)

    # Create an ElementTree object and write to the XML file
    try:
        tree = ET.ElementTree(root)
        tree.write(xml_file, encoding='utf-8', xml_declaration=True)
        print(f'Successfully created well-formed XML: {xml_file}')
    except ET.ParseError as e:
        print(f"Error while writing the XML: {e}")

def consolidate_spaces(text):
    """Consolidates multiple consecutive spaces into a single space."""
    return re.sub(r'\s+', ' ', text)

def prettify_xml(xml_file):
    try:
        # Parse the XML file
        tree = ET.parse(xml_file)
        xml_data = tree.getroot()

        # Convert the ElementTree object to a string
        xml_str = ET.tostring(xml_data, encoding='utf-8', method='xml')
        xml_str = xml_str.decode('utf-8')

        # Consolidate multiple consecutive spaces
        xml_str = consolidate_spaces(xml_str)

        # Pretty print the XML string
        xml_str = minidom.parseString(xml_str).toprettyxml(indent="  ")

        # Write the pretty XML back to the file
        with open(xml_file, 'w', encoding='utf-8') as f:
            f.write(xml_str)

        print(f'Successfully prettified and consolidated spaces in XML: {xml_file}')

    except ET.ParseError as e:
        print(f"XML parsing error: {e}")
    except FileNotFoundError:
        print(f"File not found: {xml_file}")
    except IOError as e:
        print(f"I/O error while accessing the file: {e}")
    except Exception as e:
        print(f"An unexpected error occurred: {e}")

def find_error_in_xml(xml_file):
    try:
        # Try to parse the XML file
        tree = ET.parse(xml_file)
        print(f'Successfully parsed XML: {xml_file}')
    except ET.ParseError as e:
        print(f"XML parsing error: {e}")

        # Get the line and column numbers from the error
        line_number, column_number = e.position

        # Open the file and print the specific line with the error
        with open(xml_file, 'r', encoding='utf-8') as f:
            lines = f.readlines()
            error_line = lines[line_number - 1]  # Line numbers are 1-based, so subtract 1

            print(f"Problematic line {line_number}, column {column_number}:")
            print(error_line)

            # Optionally highlight the problematic part
            print(" " * (column_number - 1) + "^")  # Print caret under the problematic character


def main():
    ppt_file_path = '../test/ML-Topic1A.pptx'  # Modify the path as needed
    xml_file_path = '../test/ML-Topic1A.xml'   # Modify the path as needed

    ppt2xml(ppt_file_path, xml_file_path)
    find_error_in_xml(xml_file_path)
    prettify_xml(xml_file_path)
    print(f'Converted {ppt_file_path} to {xml_file_path}')

if __name__ == '__main__':
    main()
