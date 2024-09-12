from pptx import Presentation
import xml.etree.ElementTree as ET
import re
from xml.sax.saxutils import escape
from xml.dom import minidom

def filter_text(text):
    # Replace control characters and invalid characters with a space or remove them
    cleaned_text = re.sub(r'[\x00-\x1F\x7F]', '', text)  # Remove control characters
    return cleaned_text

def ppt2xml(ppt_file, xml_file):
    # Load the PowerPoint presentation
    prs = Presentation(ppt_file)

    # Create the root element for the XML
    root = ET.Element('presentation')

    # Loop through each slide in the presentation
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_element = ET.SubElement(root, 'slide')

        # Get the title of the slide
        title = ''
        if slide.shapes.title:
            title = slide.shapes.title.text

        # Get the content of the slide
        content = ''
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text + '\n'

        # Filter the title and content text
        title = filter_text(title)
        content = filter_text(content)

        try:
            # Add the title and content to the slide element
            title_element = ET.SubElement(slide_element, 'title')
            title_element.text = escape(title)

            content_element = ET.SubElement(slide_element, 'content')
            content_element.text = escape(content)

        except Exception as e:
            # If there's an error, print the slide number and the problematic content
            print(f"Error processing slide {slide_num}: {e}")
            print(f"Title: {title}")
            print(f"Content: {content}")
            raise  # Re-raise the exception after logging the details

    # Create an ElementTree object and write to the XML file
    try:
        tree = ET.ElementTree(root)
        tree.write(xml_file, encoding='utf-8', xml_declaration=True)
        print(f'Successfully created well-formed XML: {xml_file}')
    except ET.ParseError as e:
        print(f"Error while writing the XML: {e}")

def prettify_xml(xml_file):
    try:
        # Parse the XML file
        tree = ET.parse(xml_file)
        xml_data = tree.getroot()

        # Convert the ElementTree object to a string
        xml_str = ET.tostring(xml_data, encoding='utf-8', method='xml')
        xml_str = xml_str.decode('utf-8')

        # Pretty print the XML string
        xml_str = minidom.parseString(xml_str).toprettyxml(indent="  ")

        # Write the pretty XML back to the file
        with open(xml_file, 'w', encoding='utf-8') as f:
            f.write(xml_str)

        print(f'Successfully prettified XML: {xml_file}')

    except ET.ParseError as e:
        print(f"XML parsing error: {e}")
    except FileNotFoundError:
        print(f"File not found: {xml_file}")
    except IOError as e:
        print(f"I/O error while accessing the file: {e}")
    except Exception as e:
        print(f"An unexpected error occurred: {e}")


def find_error_in_xml(xml_file):
    try:
        # Try to parse the XML file
        tree = ET.parse(xml_file)
        print(f'Successfully parsed XML: {xml_file}')
    except ET.ParseError as e:
        print(f"XML parsing error: {e}")

        # Get the line and column numbers from the error
        line_number, column_number = e.position

        # Open the file and print the specific line with the error
        with open(xml_file, 'r', encoding='utf-8') as f:
            lines = f.readlines()
            error_line = lines[line_number - 1]  # Line numbers are 1-based, so subtract 1

            print(f"Problematic line {line_number}, column {column_number}:")
            print(error_line)

            # Optionally highlight the problematic part
            print(" " * (column_number - 1) + "^")  # Print caret under the problematic character

def main():
    ppt_file_path = '../test/test.pptx'  # Modify the path as needed
    xml_file_path = '../test/test.xml'   # Modify the path as needed

    ppt2xml(ppt_file_path, xml_file_path)
    find_error_in_xml(xml_file_path)
    prettify_xml(xml_file_path)
    print(f'Converted {ppt_file_path} to {xml_file_path}')

if __name__ == '__main__':
    main()